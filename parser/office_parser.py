import io
import logging
from pathlib import Path
from typing import Optional, Union
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook


from striprtf.striprtf import rtf_to_text

from .office_types import (
    OfficeParserConfig, OfficeParserAST, OfficeContentNode,
    OfficeMetadata, TextFormatting, OfficeAttachment
)

logger = logging.getLogger("office_parser")

_GENERIC_TITLES = {"PowerPoint Presentation", "Word Document", "Microsoft Word Document",
                   "Microsoft PowerPoint Presentation", "Presentation", "Document"}

def _clean_title(title: Optional[str]) -> Optional[str]:
    if not title or title.strip() in _GENERIC_TITLES:
        return None
    return title.strip()



class OfficeParser:
    @staticmethod
    def parse_office(
        file: Union[str, bytes, Path],
        config: OfficeParserConfig = None
    ) -> OfficeParserAST:
        if config is None:
            config = OfficeParserConfig()

        if isinstance(file, (str, Path)):
            file_path = Path(file)
            if not file_path.exists():
                raise FileNotFoundError(f"File not found: {file}")
            ext = file_path.suffix.lower().lstrip('.')
            with open(file_path, 'rb') as f:
                data = f.read()
        elif isinstance(file, bytes):
            data = file
            ext = _detect_extension(data)
        else:
            raise ValueError("Invalid file input")

        if ext == 'docx':
            return _parse_docx(data, config)
        elif ext == 'pptx':
            return _parse_pptx(data, config)
        elif ext == 'xlsx':
            return _parse_xlsx(data, config)
        elif ext == 'rtf':
            return _parse_rtf(data, config)
        else:
            raise ValueError(f"Unsupported file type: {ext}")


def _detect_extension(data: bytes) -> str:
    if data.startswith(b'PK'):
        import zipfile
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                names = zf.namelist()
                if any(n.startswith('word/') for n in names):
                    return 'docx'
                elif any(n.startswith('ppt/') for n in names):
                    return 'pptx'
                elif any(n.startswith('xl/') for n in names):
                    return 'xlsx'
        except zipfile.BadZipFile:
            pass
        return 'docx'  # fallback
    elif data.startswith(b'{\\rtf'):
        return 'rtf'
    else:
        raise ValueError("Cannot detect file type from bytes")


def _extract_theme_colors(wb) -> list:
    """워크북에서 테마 색상 팔레트 추출"""
    try:
        import xml.etree.ElementTree as ET
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        root = ET.fromstring(wb.loaded_theme)
        colors = []
        for scheme in root.findall('.//a:clrScheme', ns):
            for child in scheme:
                for c in child:
                    val = c.get('val', c.get('lastClr', ''))
                    colors.append(val)
        return colors
    except Exception:
        return []


def _resolve_color(color_obj, theme_colors: list) -> str:
    """openpyxl Color 객체를 #RRGGBB 문자열로 변환"""
    if not color_obj:
        return None
    try:
        if color_obj.type == 'rgb' and color_obj.rgb:
            rgb = str(color_obj.rgb)
            if rgb != '00000000':
                return '#' + rgb[2:]  # AARRGGBB -> #RRGGBB
        elif color_obj.type == 'theme' and theme_colors:
            idx = color_obj.theme
            # Excel은 테마 인덱스 0<->1, 2<->3을 교차 매핑
            if idx == 0: idx = 1
            elif idx == 1: idx = 0
            elif idx == 2: idx = 3
            elif idx == 3: idx = 2
            if 0 <= idx < len(theme_colors):
                val = theme_colors[idx]
                if val.startswith('#'):
                    return val
                elif val in ('window', 'windowText'):
                    return '#FFFFFF' if val == 'window' else '#000000'
                elif len(val) == 6:
                    return '#' + val
    except Exception:
        pass
    return None


def _extract_cell_style(cell, theme_colors: list) -> dict:
    """셀의 배경색, 글자색, 볼드 정보 추출"""
    style = {}
    try:
        if cell.fill and cell.fill.fill_type == 'solid':
            bg = _resolve_color(cell.fill.fgColor, theme_colors)
            if bg:
                style['background-color'] = bg
        if cell.font:
            fc = _resolve_color(cell.font.color, theme_colors)
            if fc:
                style['color'] = fc
            if cell.font.bold:
                style['font-weight'] = 'bold'
    except Exception:
        pass
    return style if style else None


def _extract_docx_images(element, doc, qn) -> list:
    """문단/run에서 인라인 이미지 추출. 반환: [(img_bytes, ext), ...]"""
    images = []
    for blip in element.findall(f".//{qn('a:blip')}"):
        rId = blip.get(qn("r:embed"))
        if not rId:
            continue
        try:
            rel = doc.part.rels[rId]
            img_data = rel.target_part.blob
            ct = rel.target_part.content_type or ""
            ext = ct.split("/")[-1].replace("jpeg", "jpg") if "/" in ct else "png"
            images.append((img_data, ext))
        except Exception:
            pass
    return images



def _parse_docx(data: bytes, config: OfficeParserConfig) -> OfficeParserAST:
    from docx.oxml.ns import qn

    doc = Document(io.BytesIO(data))
    attachments = []
    img_counter = [0]
    # 플랫하게 모든 노드를 먼저 수집
    flat_nodes = []

    for element in doc.element.body:
        # 테이블
        if element.tag == qn("w:tbl"):
            from docx.table import Table as DocxTable
            tbl = DocxTable(element, doc)
            rows = []
            for row in tbl.rows:
                cells = [cell.text or "" for cell in row.cells]
                rows.append(OfficeContentNode(type="row", children=[
                    OfficeContentNode(type="cell", text=c) for c in cells
                ]))
            if rows:
                flat_nodes.append(OfficeContentNode(type="table", metadata={}, children=rows))
            continue

        if element.tag != qn("w:p"):
            continue

        pPr = element.find(qn("w:pPr"))

        # 이미지 추출
        if config.extract_attachments:
            for img_data, ext in _extract_docx_images(element, doc, qn):
                img_counter[0] += 1
                filename = f"image_{img_counter[0]}.{ext}"
                attachments.append(OfficeAttachment(type="image", data=img_data, filename=filename, extension=ext))
                flat_nodes.append(OfficeContentNode(
                    type="image", metadata={"filename": filename, "format": ext}
                ))

        from docx.text.paragraph import Paragraph
        para = Paragraph(element, doc)
        text = para.text.strip()
        if not text:
            continue

        style_name = para.style.name if para.style else ""

        # Heading
        if style_name.startswith("Heading"):
            try:
                level = int(style_name.split()[-1])
            except (ValueError, IndexError):
                level = 1
            flat_nodes.append(OfficeContentNode(type="heading", text=text, metadata={"level": level}))
            continue

        # List
        numPr = pPr.find(qn("w:numPr")) if pPr is not None else None
        if numPr is not None:
            ilvl_el = numPr.find(qn("w:ilvl"))
            indent_level = int(ilvl_el.get(qn("w:val"), "0")) if ilvl_el is not None else 0
            list_type = "ordered" if "Number" in style_name or "List Number" in style_name else "unordered"
            flat_nodes.append(OfficeContentNode(
                type="list", text=text, metadata={"listType": list_type, "indent_level": indent_level}
            ))
            continue

        # 일반 문단
        indent_level = 0
        if pPr is not None:
            ind = pPr.find(qn("w:ind"))
            if ind is not None:
                left = ind.get(qn("w:left"), "0")
                try:
                    indent_level = max(0, int(left) // 720)
                except ValueError:
                    pass

        flat_nodes.append(OfficeContentNode(
            type="paragraph", text=text, metadata={"indent_level": indent_level} if indent_level > 0 else None
        ))

    # -- Heading 기반 섹션 분할 --
    # 최상위 heading(가장 작은 level)을 기준으로 섹션 분할
    heading_levels = [n.metadata["level"] for n in flat_nodes if n.type == "heading" and n.metadata]
    has_sections = len(heading_levels) >= 2  # heading이 2개 이상이어야 섹션 분할 의미 있음

    if has_sections:
        split_level = min(heading_levels)
        content = []
        current_children = []
        current_title = None
        section_num = 0

        def _flush_section():
            nonlocal current_children, current_title, section_num
            if current_children:
                section_num += 1
                meta = {"sectionNumber": section_num}
                if current_title:
                    meta["sectionTitle"] = current_title
                content.append(OfficeContentNode(
                    type="section", text="", metadata=meta, children=list(current_children)
                ))
                current_children = []
                current_title = None

        for node in flat_nodes:
            if node.type == "heading" and node.metadata.get("level") == split_level:
                _flush_section()
                current_title = node.text
                current_children.append(node)
            else:
                current_children.append(node)

        _flush_section()
    else:
        # 섹션 분할 없이 플랫하게
        content = flat_nodes

    metadata = OfficeMetadata(
        title=_clean_title(doc.core_properties.title), author=doc.core_properties.author,
        created=doc.core_properties.created, modified=doc.core_properties.modified,
    )
    ast = OfficeParserAST(type="docx", metadata=metadata, content=content,
                          attachments=attachments or [])

    return ast


def _pptx_shape_to_node(shape, slide_idx: int, img_counter: list, config: OfficeParserConfig, attachments: list, slide_area: int = 0):
    """단일 shape을 OfficeContentNode로 변환"""
    from pptx.shapes.group import GroupShape
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    nodes = []

    # 테이블
    if shape.has_table:
        tbl = shape.table
        rows = []
        for row in tbl.rows:
            cells = [cell.text or "" for cell in row.cells]
            rows.append(cells)
        max_cols = max((len(r) for r in rows), default=0)
        tbl_node = OfficeContentNode(type="table", metadata={"rows": len(rows), "cols": max_cols}, children=[])
        for ri, cells in enumerate(rows):
            row_node = OfficeContentNode(type="row", metadata={"row": ri}, children=[])
            for c in cells:
                row_node.children.append(OfficeContentNode(type="cell", text=c))
            tbl_node.children.append(row_node)
        nodes.append(tbl_node)

    # 이미지 — 슬라이드 면적의 30% 미만이면 스킵
    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            if slide_area > 0:
                shape_area = shape.width * shape.height
                if shape_area / slide_area < 0.3:
                    return nodes  # 작은 이미지 — 추출/저장 모두 스킵

            img = shape.image
            ext = img.content_type.split("/")[-1].replace("jpeg", "jpg")
            img_data = img.blob
            idx = img_counter[0]
            img_counter[0] += 1
            filename = f"slide_{slide_idx}_image_{idx}.{ext}"
            meta = {"format": ext, "slideNumber": slide_idx, "imageIndex": idx, "filename": filename}
            try:
                meta["bbox"] = {
                    "left": round(shape.left / 914400, 2),
                    "top": round(shape.top / 914400, 2),
                    "width": round(shape.width / 914400, 2),
                    "height": round(shape.height / 914400, 2),
                }
            except Exception:
                pass
            img_node = OfficeContentNode(type="image", metadata=meta)
            if config.extract_attachments and img_data:
                attachments.append(OfficeAttachment(type="image", data=img_data, filename=filename, extension=ext))
            nodes.append((img_node, img_data, ext))
        except Exception:
            pass

    # 그룹 shape — 재귀
    elif isinstance(shape, GroupShape):
        for child_shape in shape.shapes:
            nodes.extend(_pptx_shape_to_node(child_shape, slide_idx, img_counter, config, attachments, slide_area))

    # 텍스트 (text_frame 있는 shape)
    elif shape.has_text_frame:
        tf = shape.text_frame
        for para in tf.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            level = para.level or 0
            p_node = OfficeContentNode(type="paragraph", text=text, metadata={"indent_level": level})
            if para.runs:
                run = para.runs[0]
                fmt = TextFormatting(bold=run.font.bold, italic=run.font.italic)
                if fmt.bold or fmt.italic:
                    p_node.formatting = fmt
            nodes.append(p_node)

    return nodes


def _parse_pptx(data: bytes, config: OfficeParserConfig) -> OfficeParserAST:
    prs = Presentation(io.BytesIO(data))
    content = []
    attachments = []

    logger.info("Parsing PowerPoint file (%d slides)", len(prs.slides))

    slide_area = prs.slide_width * prs.slide_height

    for i, slide in enumerate(prs.slides, 1):
        # 숨겨진 슬라이드 건너뛰기
        if slide._element.get('show') == '0':
            logger.info("Slide %d (hidden) — skipped", i)
            continue
        logger.info("Parsing slide %d...", i)
        slide_meta = {"slideNumber": i}

        # 슬라이드 제목 추출
        if slide.shapes.title:
            slide_meta["slideTitle"] = slide.shapes.title.text

        slide_node = OfficeContentNode(type="slide", metadata=slide_meta, children=[])
        img_counter = [0]

        for shape in slide.shapes:
            results = _pptx_shape_to_node(shape, i, img_counter, config, attachments, slide_area)
            for r in results:
                if isinstance(r, tuple):  # (img_node, img_data, ext)
                    img_node, img_data, ext = r
                    slide_node.children.append(img_node)
                else:
                    slide_node.children.append(r)

        # 슬라이드 노트
        if not config.ignore_notes and slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_node.children.append(OfficeContentNode(type="notes", text=notes_text))

        content.append(slide_node)

    metadata = OfficeMetadata(title=_clean_title(prs.core_properties.title), author=prs.core_properties.author)

    ast = OfficeParserAST(
        type="pptx", metadata=metadata, content=content,
        attachments=attachments or []
    )

    return ast



def _parse_xlsx(data: bytes, config: OfficeParserConfig) -> OfficeParserAST:
    wb = load_workbook(io.BytesIO(data), data_only=True)
    content = []
    attachments = []

    logger.info("Parsing Excel file (%d sheets)", len(wb.sheetnames))
    theme_colors = _extract_theme_colors(wb)

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        logger.info("Parsing sheet '%s'...", sheet_name)

        sheet_node = OfficeContentNode(
            type="sheet",
            metadata={"sheetName": sheet_name, "maxRow": ws.max_row, "maxColumn": ws.max_column},
            children=[]
        )
        positioned = []
        img_idx = 0

        # 이미지
        if hasattr(ws, '_images'):
            for img in ws._images:
                ext = getattr(img, 'format', 'png') or 'png'
                anchor_row, anchor_col = 0, 0
                try:
                    if hasattr(img.anchor, '_from'):
                        anchor_row = img.anchor._from.row
                        anchor_col = img.anchor._from.col
                except Exception:
                    pass

                img_meta = {"row": anchor_row, "col": anchor_col, "format": ext, "sheetName": sheet_name, "imageIndex": img_idx}
                img_node = OfficeContentNode(type="image", metadata=img_meta)
                positioned.append((anchor_row, img_node))

                img_data = None
                if config.extract_attachments:
                    try:
                        img_data = img._data()
                    except Exception:
                        pass

                # 시트별 일련번호로 파일명 생성
                safe_name = sheet_name.replace(' ', '_').replace('/', '_')
                filename = f"{safe_name}_image_{img_idx}.{ext}"
                if img_data and config.extract_attachments:
                    attachments.append(OfficeAttachment(
                        type="image", data=img_data, filename=filename, extension=ext
                    ))
                    img_node.metadata["filename"] = filename

                img_idx += 1

        # 차트
        if hasattr(ws, '_charts'):
            for chart in ws._charts:
                title = None
                if hasattr(chart, 'title') and chart.title:
                    try:
                        title = str(chart.title.tx.rich.p[0].r[0].t) if chart.title.tx and chart.title.tx.rich else None
                    except Exception:
                        pass
                chart_row = 0
                try:
                    if hasattr(chart.anchor, '_from'):
                        chart_row = chart.anchor._from.row
                except Exception:
                    pass
                chart_node = OfficeContentNode(
                    type="chart", metadata={"chartType": chart.__class__.__name__, "title": title, "row": chart_row}
                )
                positioned.append((chart_row, chart_node))

        # 병합 셀 정보 수집: (row, col) -> colspan
        merged_spans = {}
        for mc in ws.merged_cells.ranges:
            colspan = mc.max_col - mc.min_col + 1
            if colspan > 1:
                merged_spans[(mc.min_row, mc.min_col)] = colspan
                for col in range(mc.min_col + 1, mc.max_col + 1):
                    merged_spans[(mc.min_row, col)] = 0  # 병합된 나머지 셀은 스킵

        # 셀 데이터 — 빈 셀도 위치 유지를 위해 포함하되, 뒤쪽 빈 셀은 제거
        for row_idx, row in enumerate(ws.iter_rows(), 1):
            # 마지막으로 값이 있는 셀 위치 찾기 (병합 스킵 셀 제외)
            last_val_idx = -1
            for i, cell in enumerate(row):
                span = merged_spans.get((cell.row, cell.column))
                if span == 0:
                    continue  # 병합된 나머지 셀
                if cell.value is not None:
                    last_val_idx = i
            if last_val_idx < 0:
                continue
            row_node = OfficeContentNode(type="row", metadata={"row": row_idx}, children=[])
            for cell in row[:last_val_idx + 1]:
                span = merged_spans.get((cell.row, cell.column))
                if span == 0:
                    continue  # 병합된 나머지 셀 스킵
                meta = {"row": cell.row, "col": cell.column}
                if span and span > 1:
                    meta["colspan"] = span
                # 셀 스타일 수집
                style = _extract_cell_style(cell, theme_colors)
                if style:
                    meta["style"] = style
                row_node.children.append(
                    OfficeContentNode(type="cell", text=str(cell.value) if cell.value is not None else "", metadata=meta)
                )
            positioned.append((row_idx - 1, row_node))

        positioned.sort(key=lambda x: x[0])
        sheet_node.children = [node for _, node in positioned]
        content.append(sheet_node)

    return OfficeParserAST(
        type="xlsx", metadata=OfficeMetadata(), content=content,
        attachments=attachments or []
    )




def _parse_rtf(data: bytes, config: OfficeParserConfig) -> OfficeParserAST:
    text_content = rtf_to_text(data.decode('utf-8', errors='ignore'))
    content = [OfficeContentNode(type="paragraph", text=text_content)]
    return OfficeParserAST(type="rtf", metadata=OfficeMetadata(), content=content)
